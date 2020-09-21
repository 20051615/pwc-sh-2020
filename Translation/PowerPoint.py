from pptx import Presentation
from tkinter import Tk
from tkinter.filedialog import askopenfilename
from Baidu_API import translate_list, prompt_for_target_lang

if __name__ == "__main__":
    target_lang = prompt_for_target_lang()
    Tk().withdraw()

    trans_file = askopenfilename(title="Select pptx file")
    prs = Presentation(trans_file)

    text_snippets = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if len(paragraph.runs) != 0:
                        whole_paragraph = "".join(
                            run.text for run in paragraph.runs
                            )
                        text_snippets.append(whole_paragraph)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text_snippets.append(cell.text)
    
    text_snippets = translate_list(target_lang, text_snippets)
    trans_iter = iter(text_snippets)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for idx, run in enumerate(paragraph.runs):
                        if idx != 0:
                            paragraph._p.remove(run._r)
                    if len(paragraph.runs) != 0:
                        paragraph.runs[0].text = next(trans_iter)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        cell.text = next(trans_iter)

    prs.save(trans_file[:-5] + '_translated.pptx')
